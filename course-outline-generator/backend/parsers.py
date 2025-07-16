from pptx import Presentation
import fitz
from docx import Document

def parse_pptx(path, max_len=800):
    prs = Presentation(path)
    results = []
    for i, slide in enumerate(prs.slides):
        title, body = "", ""
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text.strip()
                if not title:
                    title = text
                else:
                    body += " " + text
        results.append({
            "slide": i+1,
            "title": title.strip(),
            "summary": body.strip()[:max_len]
        })
    return results

def parse_pdf(path, max_len=800):
    doc = fitz.open(path)
    results = []
    for i, page in enumerate(doc):
        lines = page.get_text().splitlines()
        title = lines[0] if lines else f"Page {i+1}"
        summary = " ".join(lines[1:])[:max_len]
        results.append({
            "slide": i+1,
            "title": title.strip(),
            "summary": summary
        })
    return results

def parse_docx(path, max_len=800):
    doc = Document(path)
    results = []
    title, body = None, ""
    count = 0
    for p in doc.paragraphs:
        if not p.text.strip():
            continue
        if p.style.name.startswith("Heading"):
            if title and body:
                results.append({
                    "slide": count+1,
                    "title": title,
                    "summary": body.strip()[:max_len]
                })
                count += 1
            title = p.text.strip()
            body = ""
        else:
            body += " " + p.text.strip()
    if title and body:
        results.append({
            "slide": count+1,
            "title": title,
            "summary": body.strip()[:max_len]
        })
    return results

def extract_modules_blocks_from_docx(path):
    doc = Document(path)
    results = []
    for table in doc.tables:
        headers = [c.text.strip().lower() for c in table.rows[0].cells]
        if "module" in headers[0] and "mvasc" in headers[1].lower():
            for row in table.rows[1:]:
                module = row.cells[0].text.strip()
                blocks = [c.text.strip() for c in row.cells[1:] if c.text.strip()]
                results.append({"module": module, "blocks": blocks})
            break
    return results
